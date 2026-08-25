import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

import content as C

GREEN = RGBColor(0x2E, 0x5D, 0x50)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
CREAM = RGBColor(0xFA, 0xF7, 0xF0)
INK = RGBColor(0x2B, 0x2B, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6D, 0x6A, 0x5E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def set_background(slide, color=CREAM):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = INK
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, color=INK, bold=False, italic=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def header_bar(slide, title, page_no):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.05), GREEN)
    add_text(slide, Inches(0.5), Inches(0.12), Inches(11.5), Inches(0.85), title,
              size=28, color=WHITE, bold=True)
    add_text(slide, Inches(12.3), Inches(7.08), Inches(0.9), Inches(0.35), str(page_no),
              size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, bullets, size=15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + b
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.font.name = "Calibri"
        p.space_after = Pt(10)
    return tb


def add_table(slide, x, y, w, h, headers, rows):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gtable = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    widths = [Inches(2.6), Inches(4.3), Inches(5.6)]
    for i, wd in enumerate(widths[:n_cols]):
        gtable.columns[i].width = wd
    for c, htext in enumerate(headers):
        cell = gtable.cell(0, c)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtable.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF7, 0xF5, 0xEF)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10.5)
                p.font.color.rgb = INK
    return gtable


# ---------------------------------------------------------------------------
# Title slide
# ---------------------------------------------------------------------------
slide = prs.slides.add_slide(BLANK)
set_background(slide, GREEN)
add_rect(slide, 0, Inches(3.0), SLIDE_W, Inches(0.05), GOLD)
add_text(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.2), C.TITLE, size=44, color=WHITE, bold=True)
add_text(slide, Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.2), C.SUBTITLE, size=17, color=CREAM, italic=True)
add_text(slide, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5),
          "github.com/dparikh066-stack/bellhaven-crm-pipeline", size=13, color=GOLD)

# ---------------------------------------------------------------------------
# Content slides
# ---------------------------------------------------------------------------
for idx, sec in enumerate(C.SECTIONS, start=2):
    slide = prs.slides.add_slide(BLANK)
    set_background(slide)
    header_bar(slide, sec["title"], idx)

    diagram = sec.get("diagram")
    is_screenshot = sec.get("screenshot", False)
    table = sec.get("table") or sec.get("big_table")
    code = sec.get("code")
    bullets = sec.get("bullets", [])

    if table:
        headers = table["headers"]
        rows = table["rows"]
        col_w = [Inches(0.7), Inches(4.3), Inches(2.1), Inches(5.2)] if len(headers) == 4 else None
        gtable = add_table(slide, Inches(0.5), Inches(1.25), Inches(12.3), Inches(5.9), headers, rows)
        if col_w:
            for i, w in enumerate(col_w):
                gtable.columns[i].width = w
    elif is_screenshot and diagram:
        img_path = os.path.join(C.IMG_ROOT, diagram)
        with Image.open(img_path) as im:
            iw, ih = im.size
        max_w, max_h = Inches(7.9), Inches(5.55)
        ratio = min(max_w / iw, max_h / ih)
        pic_w, pic_h = Emu(int(iw * ratio)), Emu(int(ih * ratio))
        slide.shapes.add_picture(img_path, Inches(0.4), Inches(1.25), width=pic_w, height=pic_h)
        add_bullets(slide, Inches(8.5), Inches(1.35), Inches(4.4), Inches(5.7), bullets, size=12.5)
    elif diagram:
        img_path = os.path.join(C.IMG_ROOT, diagram)
        with Image.open(img_path) as im:
            iw, ih = im.size
        max_w, max_h = Inches(7.6), Inches(5.7)
        ratio = min(max_w / iw, max_h / ih)
        pic_w, pic_h = Emu(int(iw * ratio)), Emu(int(ih * ratio))
        pic_x = Inches(5.4) + (max_w - pic_w) / 2
        pic_y = Inches(1.35) + (max_h - pic_h) / 2
        slide.shapes.add_picture(img_path, pic_x, pic_y, width=pic_w, height=pic_h)
        add_bullets(slide, Inches(0.5), Inches(1.35), Inches(4.7), Inches(5.7), bullets, size=13.5)
    elif code:
        add_rect(slide, Inches(0.5), Inches(1.2), Inches(7.3), Inches(5.85), RGBColor(0x1E, 0x1E, 0x1E))
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(6.9), Inches(5.55))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(code.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line if line.strip() else " "
            p.font.size = Pt(9.5)
            p.font.name = "Consolas"
            p.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
        add_bullets(slide, Inches(8.1), Inches(1.35), Inches(4.8), Inches(5.7), bullets, size=13)
    else:
        add_bullets(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.6), bullets, size=17)

out_path = os.path.join(os.path.dirname(__file__), "..", "Bellhaven_CRM_Pipeline.pptx")
prs.save(out_path)
print("Wrote", os.path.abspath(out_path))
